from django.shortcuts import get_object_or_404, redirect, render
from django.contrib.auth.views import LoginView
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.template.loader import render_to_string
import json
from decouple import config
from openai import OpenAI
from .models import ChatSession, ChatMessage
from .utils import get_gpt_response
from pptx import Presentation
import os
import markdown2
import uuid
from uuid import uuid4
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
from reportlab.lib.styles import ParagraphStyle
from django.conf import settings
from django.http import JsonResponse
from django.views.decorators.http import require_POST
from pathlib import Path
import re



api = config("OPENAI_API_KEY")
client = OpenAI(api_key=api)

BASE_DIR = Path(__file__).resolve().parent.parent

def load_system_prompt(filename):
    """Validate input."""
    path = BASE_DIR / 'system_messages' / filename
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

#===========================================================================
# PowerPoint generation with bullet points
def generate_ppt_from_text(title, bullet_points, user_id):
    """Validate input."""
    prs = Presentation()

    def set_background_color(slide, rgb_color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Research Assistant"
    set_background_color(slide, RGBColor(240, 248, 255))

    # Content Slides
    content_slide_layout = prs.slide_layouts[1]
    for point in bullet_points:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = point['title']
        set_background_color(slide, RGBColor(255, 255, 255))

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for item in point['bullets']:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = "Calibri"
            p.font.bold = False
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.space_after = Pt(16)  # Add line spacing after each bullet

    # Closing Slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Thank You!"
    slide.placeholders[1].text = ""
    set_background_color(slide, RGBColor(240, 248, 255))

    # Save file
    filename = f"ppt_{user_id}_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(settings.MEDIA_ROOT, "ppt", filename)
    prs.save(filepath)
    return settings.MEDIA_URL + f"ppt/{filename}"
#===========================================================================
# pdf generation with full academic document structure
def generate_pdf_report(json_data, user_id):
    """
    Generate PDF from GENESIS JSON schema with full academic document structure
    """
    filename = f"pdf_{user_id}_{uuid4().hex[:8]}.pdf"
    filepath = os.path.join(settings.MEDIA_ROOT, "pdf", filename)
    os.makedirs(os.path.dirname(filepath), exist_ok=True)

    doc = SimpleDocTemplate(filepath, pagesize=A4, rightMargin=50, leftMargin=50, topMargin=60, bottomMargin=60)
    styles = getSampleStyleSheet()
    story = []

    # Enhanced custom styles for academic documents
    title_style = ParagraphStyle(
        name="TitleStyle",
        fontSize=24,
        leading=30,
        alignment=TA_CENTER,
        spaceAfter=30,
        fontName='Helvetica-Bold'
    )
    
    metadata_style = ParagraphStyle(
        name="MetadataStyle",
        fontSize=11,
        leading=14,
        alignment=TA_CENTER,
        spaceAfter=8,
        textColor=colors.grey
    )
    
    heading1_style = ParagraphStyle(
        name="Heading1Style",
        fontSize=18,
        leading=22,
        spaceAfter=12,
        spaceBefore=20,
        textColor=colors.darkblue,
        fontName='Helvetica-Bold'
    )
    
    heading2_style = ParagraphStyle(
        name="Heading2Style",
        fontSize=14,
        leading=18,
        spaceAfter=10,
        spaceBefore=15,
        textColor=colors.darkblue,
        fontName='Helvetica-Bold'
    )
    
    paragraph_style = ParagraphStyle(
        name="BodyStyle",
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=12
    )
    
    citation_style = ParagraphStyle(
        name="CitationStyle",
        fontSize=10,
        leading=14,
        leftIndent=20,
        spaceAfter=6
    )

    # Extract data from JSON
    metadata = json_data.get('metadata', {})
    sections = json_data.get('sections', [])
    citations = json_data.get('citations', [])
    appendices = json_data.get('appendices', [])

    # Title Page
    story.append(Spacer(1, 1.5 * inch))
    story.append(Paragraph(metadata.get('title', 'Research Report'), title_style))
    story.append(Spacer(1, 0.5 * inch))
    
    # Metadata
    if metadata.get('author'):
        story.append(Paragraph(f"Author: {metadata['author']}", metadata_style))
    if metadata.get('date'):
        story.append(Paragraph(f"Date: {metadata['date']}", metadata_style))
    if metadata.get('document_type'):
        story.append(Paragraph(f"Document Type: {metadata['document_type']}", metadata_style))
    
    story.append(Spacer(1, 2 * inch))
    story.append(PageBreak())

    # Table of Contents (optional enhancement)
    story.append(Paragraph("Table of Contents", heading1_style))
    for idx, section in enumerate(sections, 1):
        toc_entry = f"{idx}. {section.get('heading', 'Untitled Section')}"
        story.append(Paragraph(toc_entry, styles['Normal']))
    story.append(PageBreak())

    # Main Content
    for idx, section in enumerate(sections, 1):
        section_heading = section.get('heading', 'Untitled Section')
        level = section.get('level', 1)
        
        # Choose appropriate heading style based on level
        if level == 1:
            numbered_heading = f"{idx}. {section_heading}"
            story.append(Paragraph(numbered_heading, heading1_style))
        else:
            story.append(Paragraph(section_heading, heading2_style))
        
        # Add bookmark for navigation
        story.append(Paragraph(f'<bookmark title="{section_heading}" level="{level-1}" />', styles['Normal']))
        
        # Main paragraphs
        for para in section.get("paragraphs", []):
            if para.strip():  # Only add non-empty paragraphs
                story.append(Paragraph(para, paragraph_style))
        
        # Handle subsections
        subsections = section.get('subsections', [])
        for sub_idx, subsection in enumerate(subsections, 1):
            sub_heading = subsection.get('heading', 'Untitled Subsection')
            story.append(Paragraph(f"{idx}.{sub_idx} {sub_heading}", heading2_style))
            
            for para in subsection.get("paragraphs", []):
                if para.strip():
                    story.append(Paragraph(para, paragraph_style))
        
        # Add space between major sections
        if level == 1:
            story.append(Spacer(1, 0.3 * inch))

    # References/Citations Section
    if citations:
        story.append(PageBreak())
        story.append(Paragraph("References", heading1_style))
        
        for citation in citations:
            # Format citation based on type
            citation_text = format_citation(citation)
            story.append(Paragraph(citation_text, citation_style))

    # Appendices
    if appendices:
        story.append(PageBreak())
        for app_idx, appendix in enumerate(appendices, 1):
            story.append(Paragraph(f"Appendix {app_idx}: {appendix.get('title', 'Supplementary Material')}", heading1_style))
            
            for content in appendix.get('content', []):
                if content.strip():
                    story.append(Paragraph(content, paragraph_style))

    # Build PDF with enhanced page numbering
    doc.build(story, onFirstPage=add_academic_page_number, onLaterPages=add_academic_page_number)
    return settings.MEDIA_URL + f"pdf/{filename}"


def format_citation(citation):
    """
    Format citation based on type (APA style)
    """
    authors = ", ".join(citation.get('authors', ['Unknown Author']))
    year = citation.get('year', 'n.d.')
    title = citation.get('title', 'Untitled')
    venue = citation.get('venue', '')
    
    citation_type = citation.get('type', 'journal')
    
    if citation_type == 'journal':
        return f"{authors} ({year}). {title}. <i>{venue}</i>."
    elif citation_type == 'book':
        return f"{authors} ({year}). <i>{title}</i>. {venue}."
    elif citation_type == 'conference':
        return f"{authors} ({year}). {title}. In <i>{venue}</i>."
    elif citation_type == 'web':
        return f"{authors} ({year}). {title}. Retrieved from {venue}"
    else:
        return f"{authors} ({year}). {title}. {venue}."


def add_academic_page_number(canvas, doc):
    """
    Enhanced page numbering for academic documents
    """
    page_num = canvas.getPageNumber()
    text = f"Page {page_num}"
    canvas.drawRightString(A4[0] - 50, 30, text)
    
    # Add document title in header (optional)
    canvas.setFont('Helvetica', 8)
    canvas.drawString(50, A4[1] - 30, "Research Report")



#===========================================================================
@require_POST
@login_required
def delete_chat_session(request):
    try:
        data = json.loads(request.body)
        session_id = data.get('session_id')

        session = get_object_or_404(ChatSession, id=session_id, user=request.user)
        session.delete()
        return JsonResponse({'success': True})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)

def home_view(request):
    return render(request, 'home.html')


class CustomLoginView(LoginView):
    template_name = 'authentication/login.html'
    next_page = '/'

@login_required
def chat_view(request):
    session_id = request.GET.get('session_id')
    session = None

    if session_id:
        session = ChatSession.objects.filter(id=session_id, user=request.user).first()

    messages = session.messages.all().order_by('created_at') if session else []
    sessions = ChatSession.objects.filter(user=request.user, messages__isnull=False).distinct().order_by('-created_at')

    return render(request, 'gpt/chat.html', {
        'session': session,
        'messages': messages,
        'sessions': sessions,
    })

@login_required
def get_session_list(request):
    sessions = ChatSession.objects.filter(user=request.user).order_by('-created_at')
    html = render_to_string('gpt/session_list.html', {'sessions': sessions})
    return JsonResponse({'html': html})

@login_required
def new_chat_view(request):
    return render(request, 'gpt/chat.html', {
        'messages': [],
        'session': None,
        'sessions': ChatSession.objects.filter(user=request.user).order_by('-created_at')
    })

@csrf_exempt
@login_required
def send_message(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            prompt = data.get('prompt', '')
            session_id = data.get('session_id')
            if not prompt:
                return JsonResponse({'error': 'Empty prompt'}, status=400)
            if session_id:
                session = ChatSession.objects.get(id=session_id, user=request.user)
            else:
                # Create new session only when there's a prompt to save
                session = ChatSession.objects.create(user=request.user)

            

            # session = ChatSession.objects.get(id=session_id, user=request.user)

            # Keywords
            ppt_keywords = ['ppt', 'powerpoint', 'presentation', 'slides', 'generate ppt', 'create ppt', 'make slides', 'slide deck']
            pdf_keywords = ['pdf', 'generate pdf', 'make pdf', 'export as pdf', 'create pdf', 'thesis']

            prompt_lower = prompt.lower()
            contains_pdf = any(kw in prompt_lower for kw in pdf_keywords)
            contains_ppt = any(kw in prompt_lower for kw in ppt_keywords)

            if contains_pdf or contains_ppt:
                is_pdf = contains_pdf
                is_ppt = not contains_pdf and contains_ppt

                if is_pdf:
                   system_prompt = load_system_prompt("pdf.txt")

                elif is_ppt:
                    system_prompt = load_system_prompt("ppt.txt")
                
                # Compose chat
                chat_history = [{"role": "system", "content": system_prompt}]
                for msg in session.messages.order_by('created_at'):
                    chat_history.append({"role": msg.role, "content": msg.content})
                chat_history.append({"role": "user", "content": prompt})

                # Call GPT
                # response = client.chat.completions.create(
                #     model="gpt-4o-2024-11-20",
                #     messages=chat_history,
                # )
                structured =  get_gpt_response(chat_history, return_json=True)


                # try:
                #     if not content:
                #         raise ValueError("No content returned from GPT")
                #     try:
                #         structured = json.loads(content)
                #     except json.JSONDecodeError:
                #         json_match = re.search(r'\{[\s\S]*\}', content)
                #         if json_match:
                #             json_str = json_match.group(0)
                #             structured = json.loads(json_str)
                #         else:
                #             raise ValueError("No valid JSON found in GPT content")
                # except Exception as e:
                #     return JsonResponse({'error': 'Failed to parse JSON from GPT.'}, status=500)


                if is_pdf:
                    download_url = generate_pdf_report(
                        json_data=structured,
                        user_id=request.user.id
                    )
                    file_label = "PDF Report"
                else:
                    download_url = generate_ppt_from_text(
                        title=structured['title'],
                        bullet_points=structured['slides'],
                        user_id=request.user.id
                    )
                    file_label = "PowerPoint"

                # Save messages
                ChatMessage.objects.create(session=session, role='user', content=prompt)
                ChatMessage.objects.create(session=session, role='assistant', content=f"Your {file_label} is ready: {download_url}")

                if not session.title:
                    session.title = prompt[:50].strip() + ('...' if len(prompt) > 50 else '')
                    session.save()

                return JsonResponse({'reply': f"Download your {file_label}:", 'download_url': download_url})

            # Normal chat
            default_system_prompt = load_system_prompt("default_message.txt")
            chat_history = [{"role": "system", "content":  default_system_prompt}]
            
            for msg in session.messages.order_by('created_at'):
                chat_history.append({"role": msg.role, "content": msg.content})
            chat_history.append({"role": "user", "content": prompt})

            ChatMessage.objects.create(session=session, role='user', content=prompt)

            # response = client.chat.completions.create(
            #     model="gpt-4o-2024-11-20",
            #     messages=chat_history
            # )
            reply = get_gpt_response(chat_history)

            formatted_reply = markdown2.markdown(reply)
            ChatMessage.objects.create(session=session, role='assistant', content=formatted_reply)


            if not session.title:
                session.title = prompt[:50].strip() + ('...' if len(prompt) > 50 else '')
                session.save()

            # return JsonResponse({'reply': formatted_reply})
            return JsonResponse({'reply': formatted_reply, 'session_id': str(session.id)})


        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)
